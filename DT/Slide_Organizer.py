from pptx import Presentation
from io import BytesIO

def reorganize_presentation(input_file, output_file):
    # Load the PowerPoint presentation
    prs = Presentation(input_file)
    
    # Get all slides (skip the first slide)
    slides = list(prs.slides)
    slide_count = len(slides)
    
    if slide_count <= 1:
        print("The presentation has only 1 slide or is empty. Cannot proceed.")
        return

    # Exclude the first slide and split the rest into two groups
    main_slides = slides[1:]  # Exclude the cover slide
    mid_point = len(main_slides) // 2
    
    first_half = main_slides[:mid_point]
    second_half = main_slides[mid_point:]
    
    # Create a new presentation
    new_prs = Presentation()
    new_prs.slides.add_slide(new_prs.slide_layouts[6])  # Add a blank cover slide
    
    # Interleave slides from the two groups into the new presentation
    for i in range(min(len(first_half), len(second_half))):
        _add_slide_to_presentation(new_prs, first_half[i])
        _add_slide_to_presentation(new_prs, second_half[i])
    
    # Add remaining slides (if any)
    for slide in first_half[len(second_half):] + second_half[len(first_half):]:
        _add_slide_to_presentation(new_prs, slide)
    
    # Save the new presentation
    new_prs.save(output_file)
    print(f"New presentation saved as '{output_file}'.")

def _add_slide_to_presentation(prs, slide):
    # Add a blank slide to the new presentation
    slide_layout = prs.slide_layouts[6]  # Blank slide layout
    new_slide = prs.slides.add_slide(slide_layout)
    
    # Copy all shapes and content
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        # Handle text boxes
        if shape.has_text_frame:
            new_shape = new_slide.shapes.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )
            new_text_frame = new_shape.text_frame
            for paragraph in shape.text_frame.paragraphs:
                new_paragraph = new_text_frame.add_paragraph()
                new_paragraph.text = paragraph.text
                new_paragraph.font.bold = paragraph.font.bold
                new_paragraph.font.italic = paragraph.font.italic
                new_paragraph.font.size = paragraph.font.size
        # Handle images and place them at the top
        elif shape.shape_type == 13:  # SHAPE_TYPE.PICTURE
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Calculate scaling ratio to fit image within slide width
            image_ratio = shape.width / shape.height
            new_width = slide_width
            new_height = slide_width / image_ratio
            
            # Position the image at the top
            left = 0
            top = 0
            
            # Ensure the height does not exceed the slide height
            if new_height > slide_height:
                new_height = slide_height
                new_width = slide_height * image_ratio
            
            image_stream = BytesIO(shape.image.blob)
            new_slide.shapes.add_picture(
                image_stream, left, top, int(new_width), int(new_height)
            )

# Example usage
input_pptx = "orjinal_sunum.pptx"  # Original presentation file name
output_pptx = "yeniden_siralanmis_sunum.pptx"  # Output presentation file name
reorganize_presentation(input_pptx, output_pptx)
